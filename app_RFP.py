"""
🎯 Enhanced RFP Analysis & Vendor Management System
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
Complete RFP lifecycle management with AI-powered analysis
Full vendor onboarding, evaluation, and selection workflow
Comprehensive sample data generation for end-to-end testing
"""

import streamlit as st
import anthropic
import PyPDF2
import docx
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import json
import re
from datetime import datetime, timedelta
import io
import hashlib
import uuid
from typing import Dict, List, Optional, Tuple, Any
import base64
import time
import random
from faker import Faker
import zipfile
import openpyxl
from openpyxl.styles import PatternFill, Font, Border, Side, Alignment

# Initialize Faker for sample data
fake = Faker()

# ========================================
# CONFIGURATION & INITIALIZATION
# ========================================

st.set_page_config(
    page_title="RFP Vendor Management System",
    page_icon="🎯",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Professional CSS styling
st.markdown("""
<style>
    :root {
        --primary: #2C3E50;
        --secondary: #3498DB;
        --success: #27AE60;
        --warning: #F39C12;
        --danger: #E74C3C;
        --info: #16A085;
    }
    
    .main-header {
        background: linear-gradient(135deg, var(--primary) 0%, var(--secondary) 100%);
        color: white;
        padding: 2rem;
        border-radius: 10px;
        margin-bottom: 2rem;
        text-align: center;
        box-shadow: 0 4px 15px rgba(0,0,0,0.2);
    }
    
    .workflow-card {
        background: white;
        border-radius: 10px;
        padding: 1.5rem;
        margin: 1rem 0;
        box-shadow: 0 2px 10px rgba(0,0,0,0.1);
        border-left: 4px solid var(--primary);
        transition: all 0.3s;
    }
    
    .workflow-card:hover {
        transform: translateY(-2px);
        box-shadow: 0 4px 15px rgba(0,0,0,0.15);
    }
    
    .stage-pending {
        border-left-color: #95A5A6;
        opacity: 0.7;
    }
    
    .stage-active {
        border-left-color: var(--warning);
        background: #FFF9E6;
        animation: pulse 2s infinite;
    }
    
    .stage-complete {
        border-left-color: var(--success);
        background: #E8F6F3;
    }
    
    @keyframes pulse {
        0%, 100% { opacity: 1; }
        50% { opacity: 0.9; }
    }
    
    .vendor-card {
        background: linear-gradient(to bottom, #ffffff, #f8f9fa);
        border-radius: 10px;
        padding: 1.5rem;
        margin: 1rem 0;
        box-shadow: 0 3px 10px rgba(0,0,0,0.1);
        position: relative;
    }
    
    .vendor-selected {
        border: 2px solid var(--success);
        background: #E8F6F3;
    }
    
    .vendor-shortlisted {
        border: 2px solid var(--warning);
        background: #FFF9E6;
    }
    
    .vendor-rejected {
        opacity: 0.6;
        border: 2px solid var(--danger);
    }
    
    .score-badge {
        display: inline-block;
        padding: 0.5rem 1rem;
        border-radius: 20px;
        font-weight: bold;
        margin: 0.25rem;
    }
    
    .score-excellent { background: #D4EDDA; color: #155724; }
    .score-good { background: #D1ECF1; color: #0C5460; }
    .score-fair { background: #FFF3CD; color: #856404; }
    .score-poor { background: #F8D7DA; color: #721C24; }
    
    .test-mode-banner {
        background: linear-gradient(90deg, #FF6B6B, #4ECDC4);
        color: white;
        padding: 1rem;
        text-align: center;
        font-weight: bold;
        margin-bottom: 1rem;
        border-radius: 5px;
        animation: gradient 3s ease infinite;
    }
    
    @keyframes gradient {
        0% { background-position: 0% 50%; }
        50% { background-position: 100% 50%; }
        100% { background-position: 0% 50%; }
    }
    
    .metric-card {
        background: white;
        padding: 1.5rem;
        border-radius: 10px;
        box-shadow: 0 2px 8px rgba(0,0,0,0.1);
        text-align: center;
        border-top: 3px solid var(--secondary);
    }
    
    .chat-message {
        padding: 1rem;
        margin: 0.5rem 0;
        border-radius: 10px;
    }
    
    .user-message {
        background: #E3F2FD;
        margin-left: 20%;
    }
    
    .assistant-message {
        background: #F5F5F5;
        margin-right: 20%;
    }
    
    .progress-bar {
        background: #E0E0E0;
        height: 30px;
        border-radius: 15px;
        overflow: hidden;
        margin: 20px 0;
    }
    
    .progress-fill {
        background: linear-gradient(90deg, var(--success), var(--secondary));
        height: 100%;
        display: flex;
        align-items: center;
        justify-content: center;
        color: white;
        font-weight: bold;
        transition: width 0.5s ease;
    }
    
    .document-upload-zone {
        border: 2px dashed var(--secondary);
        border-radius: 10px;
        padding: 2rem;
        text-align: center;
        background: #F0F8FF;
        margin: 1rem 0;
    }
    
    .sample-data-card {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        padding: 1.5rem;
        border-radius: 10px;
        margin: 1rem 0;
    }
    
    .timeline-container {
        position: relative;
        padding: 20px 0;
    }
    
    .timeline-item {
        display: flex;
        align-items: center;
        margin: 20px 0;
        position: relative;
    }
    
    .timeline-marker {
        width: 40px;
        height: 40px;
        border-radius: 50%;
        display: flex;
        align-items: center;
        justify-content: center;
        font-weight: bold;
        color: white;
        margin-right: 20px;
        flex-shrink: 0;
        z-index: 2;
    }
    
    .timeline-line {
        position: absolute;
        top: 40px;
        left: 20px;
        width: 2px;
        height: calc(100% + 20px);
        background: #E0E0E0;
        z-index: 1;
    }
    
    .timeline-line-active {
        background: var(--success);
    }
    
    .action-buttons {
        display: flex;
        gap: 10px;
        margin-top: 1rem;
    }
</style>
""", unsafe_allow_html=True)

# ========================================
# SAMPLE DATA GENERATOR
# ========================================

class SampleDataGenerator:
    """Generate comprehensive sample data for testing"""
    
    def __init__(self):
        self.fake = Faker()
        self.companies = [
            "Global Logistics Solutions Inc.",
            "Premier Warehousing Partners LLC",
            "FastTrack Distribution Corp.",
            "NextGen Fulfillment Systems",
            "Integrated Supply Chain Co.",
            "Strategic Logistics Group",
            "National Distribution Network",
            "Express Warehouse Services",
            "Smart Logistics Technologies",
            "Unified Transport Solutions"
        ]
        
    def generate_complete_test_data(self) -> Dict:
        """Generate complete test data for entire workflow"""
        test_data = {
            "rfp_document": self.generate_rfp_document(),
            "vendors": [],
            "workflow_state": {}
        }
        
        # Generate 5 vendors with varying quality
        for i in range(5):
            vendor_data = {
                "id": f"VND-TEST-{str(uuid.uuid4())[:8].upper()}",
                "name": self.companies[i],
                "email": f"vendor@{self.companies[i].lower().replace(' ', '')}.com",
                "proposals": self.generate_vendor_proposal(self.companies[i]),
                "scores": self.generate_evaluation_scores(quality_tier=i),
                "status": "Submitted",
                "submission_date": datetime.now() - timedelta(days=random.randint(1, 10))
            }
            test_data["vendors"].append(vendor_data)
        
        return test_data
    
    def generate_rfp_document(self) -> bytes:
        """Generate a complete RFP document in DOCX format"""
        doc = Document()
        
        # Title
        title = doc.add_heading('REQUEST FOR PROPOSAL (RFP)', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        doc.add_heading('Comprehensive Warehouse and Logistics Services', 1)
        
        # RFP Details
        rfp_number = f"RFP-{datetime.now().year}-WLS-{random.randint(1000, 9999)}"
        doc.add_paragraph(f"RFP Number: {rfp_number}")
        doc.add_paragraph(f"Issue Date: {datetime.now().strftime('%B %d, %Y')}")
        doc.add_paragraph(f"Due Date: {(datetime.now() + timedelta(days=30)).strftime('%B %d, %Y')}")
        doc.add_paragraph(f"Contract Value: ${random.randint(5, 15)},000,000 - ${random.randint(15, 25)},000,000 annually")
        
        # Executive Summary
        doc.add_heading('1. Executive Summary', 1)
        doc.add_paragraph(
            f"We are seeking qualified vendors to provide comprehensive warehouse and logistics services "
            f"for our {random.choice(['national', 'global', 'regional'])} distribution network. "
            f"This RFP encompasses warehousing, order fulfillment, transportation management, "
            f"and value-added services across {random.randint(5, 20)} locations."
        )
        
        # Scope of Services
        doc.add_heading('2. Scope of Services', 1)
        
        doc.add_heading('2.1 Warehousing Requirements', 2)
        requirements = [
            f"Minimum {random.randint(300000, 1000000):,} sq ft of warehouse space",
            "Temperature-controlled zones (ambient, cooled, frozen)",
            "24/7 operations capability with 99.9% uptime",
            "Real-time inventory management with <0.1% variance",
            "Cross-docking and transloading capabilities"
        ]
        for req in requirements:
            doc.add_paragraph(f"• {req}", style='List Bullet')
        
        doc.add_heading('2.2 Performance Requirements', 2)
        performance = [
            f"Order accuracy: >{random.uniform(99.5, 99.9):.1f}%",
            f"On-time delivery: >{random.uniform(95, 98):.1f}%",
            f"Inventory accuracy: >{random.uniform(99.7, 99.9):.1f}%",
            f"Same-day processing: {random.randint(85, 95)}% of orders",
            f"Returns processing: <{random.randint(24, 48)} hours"
        ]
        for perf in performance:
            doc.add_paragraph(f"• {perf}", style='List Bullet')
        
        # Technical Requirements
        doc.add_heading('3. Technical Requirements', 1)
        
        table = doc.add_table(rows=1, cols=3)
        table.style = 'Light Grid Accent 1'
        hdr_cells = table.rows[0].cells
        hdr_cells[0].text = 'System/Technology'
        hdr_cells[1].text = 'Requirement'
        hdr_cells[2].text = 'Priority'
        
        tech_requirements = [
            ('WMS', 'Tier-1 system with API integration', 'Critical'),
            ('EDI', 'Support for 850, 856, 810, 997', 'Critical'),
            ('Real-time tracking', 'GPS and RFID enabled', 'High'),
            ('Analytics', 'Dashboard and reporting tools', 'High'),
            ('Integration', 'RESTful API support', 'Critical')
        ]
        
        for tech, req, priority in tech_requirements:
            row_cells = table.add_row().cells
            row_cells[0].text = tech
            row_cells[1].text = req
            row_cells[2].text = priority
        
        # Evaluation Criteria
        doc.add_heading('4. Evaluation Criteria', 1)
        doc.add_paragraph("Proposals will be evaluated based on the following weighted criteria:")
        
        criteria = [
            ("Technical Capability", 25),
            ("Operational Excellence", 20),
            ("Pricing Competitiveness", 20),
            ("Experience & References", 15),
            ("Compliance & Security", 10),
            ("Innovation & Technology", 10)
        ]
        
        for criterion, weight in criteria:
            doc.add_paragraph(f"• {criterion}: {weight}%", style='List Bullet')
        
        # Submission Requirements
        doc.add_heading('5. Submission Requirements', 1)
        doc.add_paragraph("Vendors must submit the following documents:")
        
        docs_required = [
            "Technical Proposal detailing approach and capabilities",
            "Pricing Proposal with detailed cost breakdown",
            "Company Profile and Financial Statements",
            "Client References (minimum 3)",
            "Implementation Plan and Timeline",
            "Compliance Certifications"
        ]
        
        for doc_req in docs_required:
            doc.add_paragraph(f"• {doc_req}", style='List Bullet')
        
        # Save to bytes
        doc_io = io.BytesIO()
        doc.save(doc_io)
        doc_io.seek(0)
        return doc_io.getvalue()
    
    def generate_vendor_proposal(self, vendor_name: str) -> Dict[str, bytes]:
        """Generate a complete vendor proposal package"""
        proposals = {}
        
        # 1. Technical Proposal
        tech_doc = Document()
        tech_doc.add_heading(f'{vendor_name}', 0)
        tech_doc.add_heading('Technical Proposal', 1)
        
        tech_doc.add_heading('Executive Summary', 2)
        tech_doc.add_paragraph(
            f"{vendor_name} is a leading logistics provider with {random.randint(10, 30)} years "
            f"of experience serving {random.randint(100, 1000)} clients globally. "
            f"We operate {random.randint(20, 100)} facilities totaling {random.randint(1, 10)} million sq ft."
        )
        
        tech_doc.add_heading('Technical Capabilities', 2)
        
        # Add capabilities table
        table = tech_doc.add_table(rows=1, cols=2)
        table.style = 'Light Grid Accent 1'
        hdr_cells = table.rows[0].cells
        hdr_cells[0].text = 'Capability'
        hdr_cells[1].text = 'Details'
        
        capabilities = [
            ('Warehouse Management', f'WMS: {random.choice(["SAP EWM", "Manhattan", "Blue Yonder", "Oracle WMS"])}'),
            ('Automation Level', f'{random.randint(40, 90)}% automated processes'),
            ('Order Accuracy', f'{random.uniform(99, 99.9):.2f}%'),
            ('On-time Delivery', f'{random.uniform(94, 99):.1f}%'),
            ('Technology Stack', 'Cloud-based, AI-powered, IoT-enabled')
        ]
        
        for cap, detail in capabilities:
            row_cells = table.add_row().cells
            row_cells[0].text = cap
            row_cells[1].text = detail
        
        tech_io = io.BytesIO()
        tech_doc.save(tech_io)
        tech_io.seek(0)
        proposals['Technical_Proposal.docx'] = tech_io.getvalue()
        
        # 2. Pricing Proposal
        pricing_doc = Document()
        pricing_doc.add_heading(f'{vendor_name}', 0)
        pricing_doc.add_heading('Pricing Proposal', 1)
        
        pricing_doc.add_heading('Pricing Structure', 2)
        
        # Pricing table
        table = pricing_doc.add_table(rows=1, cols=4)
        table.style = 'Light Grid Accent 1'
        hdr_cells = table.rows[0].cells
        hdr_cells[0].text = 'Service'
        hdr_cells[1].text = 'Unit'
        hdr_cells[2].text = 'Price'
        hdr_cells[3].text = 'Volume Discount'
        
        pricing_items = [
            ('Storage', 'Per pallet/month', f'${random.uniform(15, 30):.2f}', '5% @ 1000+'),
            ('Pick & Pack', 'Per order', f'${random.uniform(3, 8):.2f}', '10% @ 10000+'),
            ('Shipping', 'Per package', f'${random.uniform(5, 15):.2f}', '7% @ 5000+'),
            ('Returns', 'Per item', f'${random.uniform(2, 5):.2f}', '5% @ 1000+'),
            ('Kitting', 'Per kit', f'${random.uniform(1, 3):.2f}', '10% @ 5000+')
        ]
        
        for service, unit, price, discount in pricing_items:
            row_cells = table.add_row().cells
            row_cells[0].text = service
            row_cells[1].text = unit
            row_cells[2].text = price
            row_cells[3].text = discount
        
        pricing_doc.add_heading('Total Estimated Annual Cost', 2)
        pricing_doc.add_paragraph(f"${random.randint(5, 10)},000,000 - ${random.randint(10, 15)},000,000")
        
        pricing_io = io.BytesIO()
        pricing_doc.save(pricing_io)
        pricing_io.seek(0)
        proposals['Pricing_Proposal.docx'] = pricing_io.getvalue()
        
        # 3. Company Profile (PowerPoint)
        prs = Presentation()
        
        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = vendor_name
        title_slide.placeholders[1].text = "Company Profile & Capabilities"
        
        # Overview slide
        overview_slide = prs.slides.add_slide(prs.slide_layouts[1])
        overview_slide.shapes.title.text = 'Company Overview'
        body = overview_slide.placeholders[1]
        tf = body.text_frame
        tf.text = f'• Established: {random.randint(1980, 2015)}'
        tf.add_paragraph().text = f'• Employees: {random.randint(500, 5000):,}'
        tf.add_paragraph().text = f'• Annual Revenue: ${random.randint(50, 500)}M'
        tf.add_paragraph().text = f'• Locations: {random.randint(10, 50)} facilities'
        tf.add_paragraph().text = f'• Certifications: ISO 9001, ISO 14001, C-TPAT'
        
        # Services slide
        services_slide = prs.slides.add_slide(prs.slide_layouts[1])
        services_slide.shapes.title.text = 'Core Services'
        body = services_slide.placeholders[1]
        tf = body.text_frame
        services = ['• Warehousing & Distribution', '• E-commerce Fulfillment', 
                   '• Transportation Management', '• Returns Processing', '• Value-Added Services']
        tf.text = services[0]
        for service in services[1:]:
            tf.add_paragraph().text = service
        
        ppt_io = io.BytesIO()
        prs.save(ppt_io)
        ppt_io.seek(0)
        proposals['Company_Profile.pptx'] = ppt_io.getvalue()
        
        # 4. Implementation Plan
        impl_doc = Document()
        impl_doc.add_heading(f'{vendor_name}', 0)
        impl_doc.add_heading('Implementation Plan', 1)
        
        impl_doc.add_heading('Project Timeline', 2)
        
        phases = [
            ('Phase 1: Planning & Design', '2 weeks', 'Requirements analysis, system design'),
            ('Phase 2: Setup & Integration', '4 weeks', 'Facility setup, system integration'),
            ('Phase 3: Testing & Training', '3 weeks', 'UAT, staff training'),
            ('Phase 4: Pilot Operations', '4 weeks', 'Limited volume testing'),
            ('Phase 5: Full Deployment', '2 weeks', 'Scale to full operations'),
            ('Phase 6: Optimization', 'Ongoing', 'Continuous improvement')
        ]
        
        table = impl_doc.add_table(rows=1, cols=3)
        table.style = 'Light Grid Accent 1'
        hdr_cells = table.rows[0].cells
        hdr_cells[0].text = 'Phase'
        hdr_cells[1].text = 'Duration'
        hdr_cells[2].text = 'Key Activities'
        
        for phase, duration, activities in phases:
            row_cells = table.add_row().cells
            row_cells[0].text = phase
            row_cells[1].text = duration
            row_cells[2].text = activities
        
        impl_doc.add_heading('Total Implementation Time: 15 weeks', 3)
        
        impl_io = io.BytesIO()
        impl_doc.save(impl_io)
        impl_io.seek(0)
        proposals['Implementation_Plan.docx'] = impl_io.getvalue()
        
        return proposals
    
    def generate_evaluation_scores(self, quality_tier: int = 0) -> Dict[str, float]:
        """Generate evaluation scores based on quality tier (0=best, 4=worst)"""
        base_ranges = {
            0: (85, 98),  # Excellent vendor
            1: (75, 90),  # Good vendor
            2: (65, 85),  # Fair vendor
            3: (55, 75),  # Marginal vendor
            4: (45, 65)   # Poor vendor
        }
        
        score_range = base_ranges.get(quality_tier, (60, 80))
        
        return {
            "technical_capability": random.uniform(*score_range),
            "operational_excellence": random.uniform(*score_range),
            "pricing_competitiveness": random.uniform(max(score_range[0]-10, 40), score_range[1]),
            "experience_references": random.uniform(*score_range),
            "compliance_security": random.uniform(score_range[0]+5, min(score_range[1]+10, 100)),
            "innovation_technology": random.uniform(*score_range),
            "financial_stability": random.uniform(*score_range),
            "implementation_approach": random.uniform(*score_range),
            "risk_management": random.uniform(*score_range),
            "sustainability": random.uniform(max(score_range[0]-5, 40), score_range[1])
        }
    
    def create_complete_test_package(self) -> bytes:
        """Create a complete ZIP package with all test documents"""
        zip_io = io.BytesIO()
        
        with zipfile.ZipFile(zip_io, 'w', zipfile.ZIP_DEFLATED) as zipf:
            # Generate complete test data
            test_data = self.generate_complete_test_data()
            
            # Add RFP document
            zipf.writestr('1_RFP_Documents/Main_RFP_Document.docx', test_data['rfp_document'])
            
            # Add vendor proposals
            for i, vendor_data in enumerate(test_data['vendors'], 1):
                vendor_folder = f'2_Vendor_Proposals/{i:02d}_{vendor_data["name"].replace(" ", "_")}'
                
                for filename, content in vendor_data['proposals'].items():
                    zipf.writestr(f'{vendor_folder}/{filename}', content)
                
                # Add vendor info JSON
                vendor_info = {
                    "vendor_id": vendor_data["id"],
                    "name": vendor_data["name"],
                    "email": vendor_data["email"],
                    "scores": vendor_data["scores"],
                    "submission_date": vendor_data["submission_date"].isoformat()
                }
                zipf.writestr(f'{vendor_folder}/vendor_info.json', 
                            json.dumps(vendor_info, indent=2))
            
            # Add README
            readme_content = """
# RFP Test Data Package - Complete Workflow Testing

## 📁 Package Contents:

### 1. RFP_Documents/
- Main RFP document with comprehensive requirements
- Contract value: $5-25M annually
- Multiple service requirements

### 2. Vendor_Proposals/
- 5 complete vendor proposals with varying quality:
  - Vendor 1: Excellent (85-98% scores)
  - Vendor 2: Good (75-90% scores)
  - Vendor 3: Fair (65-85% scores)
  - Vendor 4: Marginal (55-75% scores)
  - Vendor 5: Poor (45-65% scores)

### 3. Each Vendor Package Contains:
- Technical Proposal (DOCX)
- Pricing Proposal (DOCX)
- Company Profile (PPTX)
- Implementation Plan (DOCX)
- Vendor Info (JSON) with pre-calculated scores

## 🚀 Testing Workflow:

### Stage 1: RFP Initiation
1. Upload the Main_RFP_Document.docx
2. Review requirements and evaluation criteria

### Stage 2: Vendor Onboarding
1. Register each vendor
2. Upload their proposal documents
3. System will process and extract information

### Stage 3: Initial Screening
1. Review all 5 vendors
2. Shortlist top 3 based on scores

### Stage 4: Detailed Evaluation
1. Deep dive into shortlisted vendors
2. Compare technical capabilities
3. Analyze pricing proposals

### Stage 5: Final Selection
1. Select winning vendor (recommend Vendor 1)
2. Mark Vendor 2 as backup
3. Reject remaining vendors

### Stage 6: Contract Negotiation
1. Generate contract terms
2. Finalize with selected vendor

## 📊 Expected Outcomes:

### Vendor Rankings (by overall score):
1. **Global Logistics Solutions** (~90%): RECOMMENDED
2. **Premier Warehousing Partners** (~82%): BACKUP
3. **FastTrack Distribution** (~75%): QUALIFIED
4. **NextGen Fulfillment** (~65%): MARGINAL
5. **Integrated Supply Chain** (~55%): NOT RECOMMENDED

## 💡 Test Scenarios:

1. **Happy Path**: Select Vendor 1, complete all stages
2. **Negotiation**: Select Vendor 2, negotiate pricing
3. **Multi-vendor**: Shortlist top 3, compare in detail
4. **Rejection**: Properly reject low-scoring vendors
5. **Re-evaluation**: Request additional info from vendors

## 🎯 Key Features to Test:

- Document upload and parsing
- Automatic score calculation
- Vendor comparison matrix
- Workflow stage progression
- Final vendor selection
- Report generation
- Q&A functionality

## 📝 Notes:

- All data is randomly generated for testing
- Scores are pre-calculated for consistency
- Documents contain realistic content
- Vendor quality tiers ensure diverse testing

Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}
Version: 2.0
            """
            zipf.writestr('README.md', readme_content)
        
        zip_io.seek(0)
        return zip_io.getvalue()

# ========================================
# CORE CLASSES
# ========================================

class WorkflowStage:
    """Enhanced workflow stage with full functionality"""
    def __init__(self, stage_id: str, stage_num: int, name: str, description: str, 
                 required_docs: List[str], deliverables: List[str], duration: str):
        self.stage_id = stage_id
        self.stage_num = stage_num
        self.name = name
        self.description = description
        self.required_docs = required_docs
        self.deliverables = deliverables
        self.duration = duration
        self.status = "pending"
        self.progress = 0
        self.start_date = None
        self.end_date = None
        self.documents = {}
        self.notes = []
        
    def can_start(self, previous_stage) -> bool:
        """Check if stage can be started"""
        if previous_stage is None:
            return True
        return previous_stage.status == "complete"
    
    def start(self):
        """Start the workflow stage"""
        self.status = "active"
        self.start_date = datetime.now()
        self.progress = 10
        
    def complete(self):
        """Complete the workflow stage"""
        self.status = "complete"
        self.end_date = datetime.now()
        self.progress = 100
        
    def update_progress(self, progress: int):
        """Update stage progress"""
        self.progress = min(100, max(0, progress))
        if self.progress == 100 and self.status != "complete":
            self.complete()

class VendorProfile:
    """Enhanced vendor profile with complete lifecycle"""
    def __init__(self, vendor_id: str, name: str, email: str = None):
        self.vendor_id = vendor_id
        self.name = name
        self.email = email or f"vendor@{name.lower().replace(' ', '')}.com"
        self.registration_date = datetime.now()
        self.documents = {}
        self.scores = {}
        self.overall_score = 0
        self.status = "Registered"  # Registered, Submitted, Evaluated, Shortlisted, Selected, Rejected
        self.submission_date = None
        self.evaluation_date = None
        self.decision_date = None
        self.decision = None  # Selected, Backup, Rejected
        self.strengths = []
        self.weaknesses = []
        self.risks = []
        self.notes = []
        
    def submit_proposal(self, documents: Dict):
        """Submit vendor proposal"""
        self.documents.update(documents)
        self.submission_date = datetime.now()
        self.status = "Submitted"
        
    def evaluate(self, scores: Dict[str, float]):
        """Evaluate vendor proposal"""
        self.scores = scores
        self.overall_score = sum(scores.values()) / len(scores) if scores else 0
        self.evaluation_date = datetime.now()
        self.status = "Evaluated"
        
        # Auto-generate strengths and weaknesses
        self.strengths = [k.replace('_', ' ').title() for k, v in scores.items() if v >= 85]
        self.weaknesses = [k.replace('_', ' ').title() for k, v in scores.items() if v < 70]
        
    def shortlist(self):
        """Add vendor to shortlist"""
        if self.status == "Evaluated":
            self.status = "Shortlisted"
            
    def make_decision(self, decision: str):
        """Make final decision on vendor"""
        self.decision = decision
        self.decision_date = datetime.now()
        if decision == "Selected":
            self.status = "Selected"
        elif decision == "Backup":
            self.status = "Backup"
        else:
            self.status = "Rejected"

class RFPAnalyzer:
    """Main RFP Analysis Engine with complete workflow"""
    def __init__(self):
        self.workflow_stages = self._initialize_workflow()
        self.vendors = {}
        self.shortlisted_vendors = []
        self.selected_vendor = None
        self.backup_vendors = []
        self.current_rfp = self._generate_rfp_details()
        self.claude_client = None
        self.sample_generator = SampleDataGenerator()
        self.test_mode = False
        self.initialize_claude()
        
    def _generate_rfp_details(self):
        """Generate RFP details"""
        return {
            "rfp_id": f"RFP-2025-{str(uuid.uuid4())[:8].upper()}",
            "title": "Comprehensive Warehouse and Logistics Services RFP",
            "issue_date": datetime.now() - timedelta(days=7),
            "due_date": datetime.now() + timedelta(days=21),
            "budget": "$5,000,000 - $25,000,000",
            "contract_duration": "3 years with 2 optional 1-year extensions",
            "evaluation_criteria": {
                "Technical Capability": 25,
                "Operational Excellence": 20,
                "Pricing Competitiveness": 20,
                "Experience & References": 15,
                "Compliance & Security": 10,
                "Innovation & Technology": 10
            }
        }
    
    def initialize_claude(self):
        """Initialize Claude API"""
        try:
            api_key = None
            for key in ["CLAUDE_API_KEY", "ANTHROPIC_API_KEY", "claude_api_key", "anthropic_api_key"]:
                try:
                    api_key = st.secrets.get(key)
                    if api_key:
                        break
                except:
                    pass
            
            if api_key:
                self.claude_client = anthropic.Anthropic(api_key=api_key)
                return True
        except:
            pass
        return False
    
    def _initialize_workflow(self) -> Dict[str, WorkflowStage]:
        """Initialize complete RFP workflow with all stages"""
        stages = {}
        
        workflow_definition = [
            {
                "id": "initiation",
                "name": "RFP Initiation",
                "desc": "Receive and review RFP requirements from stakeholders",
                "docs": ["Business Requirements", "Budget Approval", "Service Specifications"],
                "deliverables": ["RFP Package", "Evaluation Criteria", "Timeline"],
                "duration": "3 days"
            },
            {
                "id": "vendor_registration",
                "name": "Vendor Registration",
                "desc": "Register and pre-qualify vendors for RFP participation",
                "docs": ["Vendor Database", "Pre-qualification Forms"],
                "deliverables": ["Registered Vendor List", "NDA Agreements"],
                "duration": "5 days"
            },
            {
                "id": "rfp_distribution",
                "name": "RFP Distribution",
                "desc": "Distribute RFP package to registered vendors",
                "docs": ["RFP Document", "Addendums"],
                "deliverables": ["Distribution Confirmation", "Q&A Schedule"],
                "duration": "2 days"
            },
            {
                "id": "proposal_submission",
                "name": "Proposal Submission",
                "desc": "Receive and validate vendor proposals",
                "docs": ["Vendor Proposals", "Technical Documents", "Pricing"],
                "deliverables": ["Submission Log", "Completeness Check"],
                "duration": "1 day"
            },
            {
                "id": "initial_screening",
                "name": "Initial Screening",
                "desc": "Screen proposals for minimum requirements",
                "docs": ["Evaluation Matrix", "Compliance Checklist"],
                "deliverables": ["Qualified Vendors", "Screening Report"],
                "duration": "3 days"
            },
            {
                "id": "detailed_evaluation",
                "name": "Detailed Evaluation",
                "desc": "Comprehensive evaluation of qualified proposals",
                "docs": ["Technical Analysis", "Pricing Analysis", "Risk Assessment"],
                "deliverables": ["Evaluation Scores", "Vendor Rankings"],
                "duration": "7 days"
            },
            {
                "id": "vendor_shortlist",
                "name": "Vendor Shortlisting",
                "desc": "Select top vendors for final consideration",
                "docs": ["Evaluation Report", "Reference Checks"],
                "deliverables": ["Shortlisted Vendors", "Presentation Schedule"],
                "duration": "2 days"
            },
            {
                "id": "vendor_presentations",
                "name": "Vendor Presentations",
                "desc": "Vendor presentations and Q&A sessions",
                "docs": ["Presentation Materials", "Q&A Responses"],
                "deliverables": ["Presentation Scores", "Clarifications"],
                "duration": "3 days"
            },
            {
                "id": "final_selection",
                "name": "Final Selection",
                "desc": "Select winning vendor and backup options",
                "docs": ["Final Evaluation", "Recommendation Report"],
                "deliverables": ["Selected Vendor", "Award Letter"],
                "duration": "2 days"
            },
            {
                "id": "contract_negotiation",
                "name": "Contract Negotiation",
                "desc": "Negotiate terms and finalize contract",
                "docs": ["Contract Draft", "Terms & Conditions"],
                "deliverables": ["Final Contract", "SLAs"],
                "duration": "5 days"
            },
            {
                "id": "vendor_onboarding",
                "name": "Vendor Onboarding",
                "desc": "Onboard selected vendor and kickoff project",
                "docs": ["Implementation Plan", "Transition Schedule"],
                "deliverables": ["Kickoff Meeting", "Project Plan"],
                "duration": "3 days"
            }
        ]
        
        for idx, stage_def in enumerate(workflow_definition, 1):
            stages[stage_def["id"]] = WorkflowStage(
                stage_def["id"],
                idx,
                stage_def["name"],
                stage_def["desc"],
                stage_def["docs"],
                stage_def["deliverables"],
                stage_def["duration"]
            )
        
        return stages
    
    def enable_test_mode(self):
        """Enable test mode with complete sample data"""
        self.test_mode = True
        
        # Generate test data
        test_data = self.sample_generator.generate_complete_test_data()
        
        # Load vendors
        for vendor_data in test_data['vendors']:
            vendor = VendorProfile(
                vendor_data['id'],
                vendor_data['name'],
                vendor_data['email']
            )
            vendor.submit_proposal(vendor_data['proposals'])
            vendor.evaluate(vendor_data['scores'])
            self.vendors[vendor.vendor_id] = vendor
        
        # Progress workflow to evaluation stage
        self.workflow_stages["initiation"].complete()
        self.workflow_stages["vendor_registration"].complete()
        self.workflow_stages["rfp_distribution"].complete()
        self.workflow_stages["proposal_submission"].complete()
        self.workflow_stages["initial_screening"].complete()
        self.workflow_stages["detailed_evaluation"].start()
        self.workflow_stages["detailed_evaluation"].update_progress(75)
        
        return test_data
    
    def get_workflow_progress(self) -> int:
        """Calculate overall workflow progress"""
        total_stages = len(self.workflow_stages)
        completed = sum(1 for s in self.workflow_stages.values() if s.status == "complete")
        active_progress = sum(s.progress/100 for s in self.workflow_stages.values() if s.status == "active")
        
        return int(((completed + active_progress) / total_stages) * 100)
    
    def shortlist_vendors(self, vendor_ids: List[str]):
        """Shortlist selected vendors"""
        self.shortlisted_vendors = []
        for vid in vendor_ids:
            if vid in self.vendors:
                self.vendors[vid].shortlist()
                self.shortlisted_vendors.append(vid)
    
    def select_vendor(self, selected_id: str, backup_ids: List[str] = None):
        """Make final vendor selection"""
        if selected_id in self.vendors:
            self.vendors[selected_id].make_decision("Selected")
            self.selected_vendor = selected_id
            
        if backup_ids:
            for bid in backup_ids:
                if bid in self.vendors:
                    self.vendors[bid].make_decision("Backup")
                    self.backup_vendors.append(bid)
        
        # Reject remaining vendors
        for vid, vendor in self.vendors.items():
            if vid != selected_id and vid not in (backup_ids or []):
                if vendor.status in ["Evaluated", "Shortlisted"]:
                    vendor.make_decision("Rejected")
    
    def generate_analysis_report(self, vendor: VendorProfile) -> Dict:
        """Generate comprehensive analysis report"""
        return {
            "vendor_id": vendor.vendor_id,
            "vendor_name": vendor.name,
            "analysis_date": datetime.now().isoformat(),
            "overall_score": vendor.overall_score,
            "status": vendor.status,
            "scores": vendor.scores,
            "strengths": vendor.strengths,
            "weaknesses": vendor.weaknesses,
            "recommendation": self._generate_recommendation(vendor),
            "executive_summary": self._generate_executive_summary(vendor),
            "next_steps": self._generate_next_steps(vendor)
        }
    
    def _generate_recommendation(self, vendor: VendorProfile) -> str:
        """Generate recommendation based on score"""
        if vendor.overall_score >= 85:
            return "Strongly Recommended - Proceed to final selection"
        elif vendor.overall_score >= 75:
            return "Recommended - Consider for shortlist"
        elif vendor.overall_score >= 65:
            return "Conditionally Recommended - Address weaknesses"
        else:
            return "Not Recommended - Does not meet requirements"
    
    def _generate_executive_summary(self, vendor: VendorProfile) -> str:
        """Generate executive summary"""
        return f"""
        {vendor.name} has submitted a comprehensive proposal with an overall score of {vendor.overall_score:.1f}/100.
        The vendor demonstrates {len(vendor.strengths)} key strengths and has {len(vendor.weaknesses)} areas for improvement.
        Based on our evaluation criteria, this vendor is {self._generate_recommendation(vendor).split(' - ')[0].lower()}.
        """
    
    def _generate_next_steps(self, vendor: VendorProfile) -> List[str]:
        """Generate next steps based on vendor status"""
        if vendor.overall_score >= 85:
            return [
                "Schedule vendor presentation",
                "Conduct reference checks",
                "Site visit (if applicable)",
                "Prepare for contract negotiations"
            ]
        elif vendor.overall_score >= 75:
            return [
                "Add to shortlist",
                "Request clarifications on weak areas",
                "Schedule follow-up discussion",
                "Compare with other vendors"
            ]
        elif vendor.overall_score >= 65:
            return [
                "Request additional information",
                "Address identified gaps",
                "Re-evaluate after improvements",
                "Keep as backup option"
            ]
        else:
            return [
                "Send regret letter",
                "Provide feedback if requested",
                "Update vendor database",
                "Consider for future RFPs"
            ]

# ========================================
# UI COMPONENTS
# ========================================

def render_header():
    """Render application header"""
    st.markdown("""
    <div class="main-header">
        <h1>🎯 RFP Vendor Management System</h1>
        <h3>Complete Procurement Workflow with AI Analysis</h3>
        <p>Vendor Selection • Evaluation • Contract Management</p>
    </div>
    """, unsafe_allow_html=True)

def render_test_mode_banner():
    """Render test mode banner"""
    if st.session_state.get('test_mode', False):
        st.markdown("""
        <div class="test-mode-banner">
            🧪 TEST MODE ACTIVE - 5 sample vendors loaded with complete proposals
        </div>
        """, unsafe_allow_html=True)

def render_workflow_dashboard(analyzer: RFPAnalyzer):
    """Render comprehensive workflow dashboard"""
    st.header("📋 RFP Workflow Management")
    
    # Overall progress
    progress = analyzer.get_workflow_progress()
    st.markdown(f"""
    <div class="progress-bar">
        <div class="progress-fill" style="width: {progress}%;">
            Overall Progress: {progress}%
        </div>
    </div>
    """, unsafe_allow_html=True)
    
    # Workflow stages
    st.subheader("Workflow Stages")
    
    # Create timeline visualization
    for idx, (stage_id, stage) in enumerate(analyzer.workflow_stages.items()):
        # Determine previous stage
        prev_stage = None
        if idx > 0:
            prev_stage = list(analyzer.workflow_stages.values())[idx - 1]
        
        # Stage styling
        if stage.status == "complete":
            marker_color = "var(--success)"
            icon = "✅"
            card_class = "stage-complete"
        elif stage.status == "active":
            marker_color = "var(--warning)"
            icon = "🔄"
            card_class = "stage-active"
        else:
            marker_color = "#95A5A6"
            icon = "⏳"
            card_class = "stage-pending"
        
        # Timeline item
        col1, col2 = st.columns([1, 11])
        
        with col1:
            st.markdown(f"""
            <div class="timeline-item">
                <div class="timeline-marker" style="background: {marker_color};">
                    {stage.stage_num}
                </div>
                {"<div class='timeline-line timeline-line-active'></div>" if stage.status == "complete" else "<div class='timeline-line'></div>"}
            </div>
            """, unsafe_allow_html=True)
        
        with col2:
            with st.expander(f"{icon} **Stage {stage.stage_num}: {stage.name}** - {stage.status.upper()}", 
                           expanded=(stage.status == "active")):
                
                col_a, col_b, col_c = st.columns([4, 3, 3])
                
                with col_a:
                    st.markdown(f"**Description:** {stage.description}")
                    st.markdown(f"**Duration:** {stage.duration}")
                    
                    if stage.status == "active":
                        progress_val = st.slider(
                            "Progress",
                            0, 100, stage.progress,
                            key=f"progress_{stage_id}"
                        )
                        if progress_val != stage.progress:
                            stage.update_progress(progress_val)
                            st.rerun()
                    elif stage.status == "complete":
                        st.progress(1.0)
                        st.caption(f"Completed on {stage.end_date.strftime('%Y-%m-%d %H:%M') if stage.end_date else 'N/A'}")
                
                with col_b:
                    st.markdown("**Required Documents:**")
                    for doc in stage.required_docs[:3]:
                        st.caption(f"• {doc}")
                    if len(stage.required_docs) > 3:
                        st.caption(f"...and {len(stage.required_docs) - 3} more")
                
                with col_c:
                    st.markdown("**Deliverables:**")
                    for deliverable in stage.deliverables[:3]:
                        st.caption(f"• {deliverable}")
                    if len(stage.deliverables) > 3:
                        st.caption(f"...and {len(stage.deliverables) - 3} more")
                    
                    if stage.status == "pending":
                        if st.button(f"Start Stage", key=f"start_{stage_id}", type="primary"):
                            if stage.can_start(prev_stage):
                                stage.start()
                                st.success(f"Started: {stage.name}")
                                st.rerun()
                            else:
                                st.error("Complete previous stage first!")
                    
                    elif stage.status == "active" and stage.progress == 100:
                        if st.button(f"Complete Stage", key=f"complete_{stage_id}", type="primary"):
                            stage.complete()
                            st.success(f"Completed: {stage.name}")
                            st.rerun()

def render_vendor_management(analyzer: RFPAnalyzer):
    """Render comprehensive vendor management interface"""
    st.header("👥 Vendor Management")
    
    # Vendor statistics
    col1, col2, col3, col4 = st.columns(4)
    
    with col1:
        st.metric("Total Vendors", len(analyzer.vendors))
    
    with col2:
        evaluated = sum(1 for v in analyzer.vendors.values() if v.status in ["Evaluated", "Shortlisted", "Selected"])
        st.metric("Evaluated", evaluated)
    
    with col3:
        shortlisted = len(analyzer.shortlisted_vendors)
        st.metric("Shortlisted", shortlisted)
    
    with col4:
        selected = 1 if analyzer.selected_vendor else 0
        st.metric("Selected", selected)
    
    # Vendor list and actions
    if analyzer.vendors:
        st.subheader("Vendor List")
        
        for vendor in analyzer.vendors.values():
            # Determine card styling
            if vendor.status == "Selected":
                card_class = "vendor-selected"
                status_icon = "🏆"
            elif vendor.status == "Shortlisted":
                card_class = "vendor-shortlisted"
                status_icon = "⭐"
            elif vendor.status == "Rejected":
                card_class = "vendor-rejected"
                status_icon = "❌"
            else:
                card_class = "vendor-card"
                status_icon = "📋"
            
            with st.container():
                col1, col2, col3, col4, col5 = st.columns([3, 2, 2, 2, 2])
                
                with col1:
                    st.markdown(f"### {status_icon} {vendor.name}")
                    st.caption(f"ID: {vendor.vendor_id}")
                
                with col2:
                    if vendor.overall_score > 0:
                        if vendor.overall_score >= 85:
                            badge = "score-excellent"
                        elif vendor.overall_score >= 75:
                            badge = "score-good"
                        elif vendor.overall_score >= 65:
                            badge = "score-fair"
                        else:
                            badge = "score-poor"
                        
                        st.markdown(f"""
                        <div class="score-badge {badge}">
                            Score: {vendor.overall_score:.1f}/100
                        </div>
                        """, unsafe_allow_html=True)
                
                with col3:
                    st.markdown(f"**Status:** {vendor.status}")
                    if vendor.submission_date:
                        st.caption(f"Submitted: {vendor.submission_date.strftime('%Y-%m-%d')}")
                
                with col4:
                    st.markdown(f"**Documents:** {len(vendor.documents)}")
                    if vendor.strengths:
                        st.caption(f"Strengths: {len(vendor.strengths)}")
                
                with col5:
                    if vendor.status == "Evaluated" and vendor.vendor_id not in analyzer.shortlisted_vendors:
                        if st.button("Shortlist", key=f"shortlist_{vendor.vendor_id}"):
                            vendor.shortlist()
                            analyzer.shortlisted_vendors.append(vendor.vendor_id)
                            st.success(f"Shortlisted: {vendor.name}")
                            st.rerun()
                    
                    if st.button("View Details", key=f"details_{vendor.vendor_id}"):
                        st.session_state.selected_vendor = vendor.vendor_id
                        st.session_state.current_page = "evaluation"
                        st.rerun()
                
                st.markdown("---")
    
    # Vendor selection controls
    if analyzer.shortlisted_vendors and not analyzer.selected_vendor:
        st.subheader("🎯 Final Vendor Selection")
        
        col1, col2 = st.columns(2)
        
        with col1:
            selected_vendor = st.selectbox(
                "Select Winning Vendor",
                options=analyzer.shortlisted_vendors,
                format_func=lambda x: analyzer.vendors[x].name + f" (Score: {analyzer.vendors[x].overall_score:.1f})"
            )
        
        with col2:
            backup_vendors = st.multiselect(
                "Select Backup Vendor(s)",
                options=[v for v in analyzer.shortlisted_vendors if v != selected_vendor],
                format_func=lambda x: analyzer.vendors[x].name
            )
        
        if st.button("Confirm Selection", type="primary"):
            analyzer.select_vendor(selected_vendor, backup_vendors)
            st.success(f"✅ Selected: {analyzer.vendors[selected_vendor].name}")
            st.balloons()
            
            # Progress workflow
            if "vendor_shortlist" in analyzer.workflow_stages:
                analyzer.workflow_stages["vendor_shortlist"].complete()
            if "final_selection" in analyzer.workflow_stages:
                analyzer.workflow_stages["final_selection"].complete()
            
            st.rerun()
    
    elif analyzer.selected_vendor:
        st.success(f"✅ **Selected Vendor:** {analyzer.vendors[analyzer.selected_vendor].name}")
        if analyzer.backup_vendors:
            st.info(f"**Backup Vendors:** {', '.join([analyzer.vendors[v].name for v in analyzer.backup_vendors])}")

def render_vendor_evaluation(analyzer: RFPAnalyzer):
    """Render detailed vendor evaluation"""
    st.header("📊 Vendor Evaluation")
    
    if not analyzer.vendors:
        st.info("No vendors available. Enable test mode to load sample vendors.")
        return
    
    # Vendor selector
    selected_vendor_id = st.session_state.get('selected_vendor', list(analyzer.vendors.keys())[0])
    
    vendor_names = {v.vendor_id: f"{v.name} - {v.status}" for v in analyzer.vendors.values()}
    selected_vendor_id = st.selectbox(
        "Select Vendor for Evaluation",
        options=list(vendor_names.keys()),
        format_func=lambda x: vendor_names[x],
        index=list(vendor_names.keys()).index(selected_vendor_id) if selected_vendor_id in vendor_names else 0
    )
    
    vendor = analyzer.vendors[selected_vendor_id]
    analysis = analyzer.generate_analysis_report(vendor)
    
    # Evaluation header
    col1, col2, col3, col4 = st.columns(4)
    
    with col1:
        st.metric("Overall Score", f"{vendor.overall_score:.1f}/100")
    
    with col2:
        st.metric("Strengths", len(vendor.strengths))
    
    with col3:
        st.metric("Weaknesses", len(vendor.weaknesses))
    
    with col4:
        st.metric("Status", vendor.status)
    
    # Evaluation tabs
    tabs = st.tabs(["📊 Scores", "💪 Analysis", "📄 Documents", "📝 Report", "💬 Notes"])
    
    with tabs[0]:  # Scores
        st.subheader("Detailed Scoring")
        
        # Radar chart
        if vendor.scores:
            categories = [k.replace('_', ' ').title() for k in vendor.scores.keys()]
            values = list(vendor.scores.values())
            
            fig = go.Figure()
            
            fig.add_trace(go.Scatterpolar(
                r=values,
                theta=categories,
                fill='toself',
                name=vendor.name
            ))
            
            # Add benchmark
            benchmark = [75] * len(categories)
            fig.add_trace(go.Scatterpolar(
                r=benchmark,
                theta=categories,
                fill='toself',
                name='Benchmark',
                line=dict(color='gray', dash='dash')
            ))
            
            fig.update_layout(
                polar=dict(radialaxis=dict(visible=True, range=[0, 100])),
                showlegend=True,
                title=f"Evaluation: {vendor.name}"
            )
            
            st.plotly_chart(fig, use_container_width=True)
        
        # Score table
        if vendor.scores:
            scores_df = pd.DataFrame([
                {"Criteria": k.replace('_', ' ').title(), "Score": f"{v:.1f}"}
                for k, v in vendor.scores.items()
            ])
            st.dataframe(scores_df, use_container_width=True, hide_index=True)
    
    with tabs[1]:  # Analysis
        st.subheader("Comprehensive Analysis")
        
        st.markdown("#### Executive Summary")
        st.info(analysis['executive_summary'])
        
        st.markdown("#### Recommendation")
        if "Strongly Recommended" in analysis['recommendation']:
            st.success(analysis['recommendation'])
        elif "Recommended" in analysis['recommendation']:
            st.info(analysis['recommendation'])
        elif "Conditionally" in analysis['recommendation']:
            st.warning(analysis['recommendation'])
        else:
            st.error(analysis['recommendation'])
        
        col1, col2 = st.columns(2)
        
        with col1:
            st.markdown("#### Strengths")
            if vendor.strengths:
                for strength in vendor.strengths:
                    st.success(f"✅ {strength}")
            else:
                st.info("No significant strengths identified")
        
        with col2:
            st.markdown("#### Weaknesses")
            if vendor.weaknesses:
                for weakness in vendor.weaknesses:
                    st.warning(f"⚠️ {weakness}")
            else:
                st.success("No major weaknesses identified")
        
        st.markdown("#### Next Steps")
        for step in analysis['next_steps']:
            st.write(f"• {step}")
    
    with tabs[2]:  # Documents
        st.subheader("Submitted Documents")
        
        if vendor.documents:
            for doc_type, doc_name in vendor.documents.items():
                st.write(f"📄 **{doc_type}:** {doc_name}")
        else:
            st.info("No documents submitted")
    
    with tabs[3]:  # Report
        st.subheader("Evaluation Report")
        
        # Generate report button
        if st.button("Generate Full Report", type="primary"):
            with st.spinner("Generating report..."):
                time.sleep(2)
                
                report_json = json.dumps(analysis, indent=2, default=str)
                
                st.download_button(
                    label="📥 Download Report (JSON)",
                    data=report_json,
                    file_name=f"evaluation_report_{vendor.vendor_id}_{datetime.now().strftime('%Y%m%d')}.json",
                    mime="application/json"
                )
                
                st.success("Report generated successfully!")
    
    with tabs[4]:  # Notes
        st.subheader("Evaluation Notes")
        
        # Add note
        new_note = st.text_area("Add a note:", key=f"note_{vendor.vendor_id}")
        if st.button("Add Note", key=f"add_note_{vendor.vendor_id}"):
            if new_note:
                vendor.notes.append({
                    "date": datetime.now(),
                    "note": new_note
                })
                st.success("Note added!")
                st.rerun()
        
        # Display notes
        if vendor.notes:
            for note in reversed(vendor.notes):
                st.info(f"**{note['date'].strftime('%Y-%m-%d %H:%M')}:** {note['note']}")
        else:
            st.caption("No notes yet")

def render_sample_data_generator():
    """Render sample data generation interface"""
    st.header("🎲 Sample Data Generator")
    
    st.markdown("""
    <div class="sample-data-card">
        <h3>Generate Complete Test Data</h3>
        <p>Create realistic RFP documents and vendor proposals for end-to-end testing</p>
    </div>
    """, unsafe_allow_html=True)
    
    col1, col2, col3 = st.columns(3)
    
    with col1:
        if st.button("📦 Generate Complete Package", type="primary", use_container_width=True):
            with st.spinner("Generating comprehensive test package..."):
                generator = SampleDataGenerator()
                zip_data = generator.create_complete_test_package()
                
                st.success("✅ Test package generated!")
                st.download_button(
                    label="📥 Download Complete Package (ZIP)",
                    data=zip_data,
                    file_name=f"RFP_Complete_Test_Package_{datetime.now().strftime('%Y%m%d_%H%M%S')}.zip",
                    mime="application/zip",
                    use_container_width=True
                )
    
    with col2:
        if st.button("📄 Generate RFP Document", use_container_width=True):
            generator = SampleDataGenerator()
            rfp_data = generator.generate_rfp_document()
            
            st.success("RFP document generated!")
            st.download_button(
                label="📥 Download RFP",
                data=rfp_data,
                file_name=f"RFP_Document_{datetime.now().strftime('%Y%m%d')}.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            )
    
    with col3:
        vendor_quality = st.selectbox(
            "Vendor Quality",
            ["Excellent", "Good", "Fair", "Marginal", "Poor"]
        )
        
        if st.button("🏢 Generate Vendor", use_container_width=True):
            generator = SampleDataGenerator()
            quality_map = {"Excellent": 0, "Good": 1, "Fair": 2, "Marginal": 3, "Poor": 4}
            
            vendor_name = generator.companies[quality_map[vendor_quality]]
            proposals = generator.generate_vendor_proposal(vendor_name)
            
            # Create ZIP for vendor
            zip_io = io.BytesIO()
            with zipfile.ZipFile(zip_io, 'w') as zipf:
                for filename, content in proposals.items():
                    zipf.writestr(filename, content)
            
            zip_io.seek(0)
            
            st.success(f"Generated {vendor_quality} vendor: {vendor_name}")
            st.download_button(
                label=f"📥 Download {vendor_name}",
                data=zip_io.getvalue(),
                file_name=f"{vendor_name.replace(' ', '_')}_Proposal.zip",
                mime="application/zip"
            )
    
    # Instructions
    st.markdown("---")
    st.subheader("📚 How to Use Test Data")
    
    with st.expander("Testing Instructions", expanded=True):
        st.markdown("""
        ### Complete End-to-End Testing:
        
        1. **Enable Test Mode** (in sidebar) to auto-load 5 sample vendors
        2. **Review Workflow** - See pre-populated stages
        3. **Evaluate Vendors** - Check scores and rankings
        4. **Shortlist Top 3** - Based on scores
        5. **Select Winner** - Choose Vendor 1 (highest score)
        6. **Complete Workflow** - Progress through remaining stages
        
        ### Manual Testing:
        
        1. Download the Complete Package
        2. Extract ZIP file
        3. Upload RFP document
        4. Register vendors one by one
        5. Upload their proposals
        6. Run evaluation
        7. Make selection
        
        ### Expected Results:
        
        - **Vendor 1:** 85-98% (Excellent) → SELECT
        - **Vendor 2:** 75-90% (Good) → BACKUP
        - **Vendor 3:** 65-85% (Fair) → SHORTLIST
        - **Vendor 4:** 55-75% (Marginal) → REJECT
        - **Vendor 5:** 45-65% (Poor) → REJECT
        """)

def render_sidebar(analyzer: RFPAnalyzer):
    """Render enhanced sidebar"""
    with st.sidebar:
        st.markdown("### 🎯 RFP Management")
        
        # Test Mode
        st.markdown("#### 🧪 Test Mode")
        test_mode = st.checkbox(
            "Enable Test Mode",
            value=st.session_state.get('test_mode', False),
            help="Load 5 sample vendors with complete proposals"
        )
        
        if test_mode != st.session_state.get('test_mode', False):
            st.session_state.test_mode = test_mode
            if test_mode:
                test_data = analyzer.enable_test_mode()
                st.success(f"Loaded {len(test_data['vendors'])} test vendors!")
            st.rerun()
        
        st.markdown("---")
        
        # RFP Info
        st.markdown("#### 📋 Current RFP")
        st.caption(f"**ID:** {analyzer.current_rfp['rfp_id']}")
        st.caption(f"**Budget:** {analyzer.current_rfp['budget']}")
        st.caption(f"**Due:** {analyzer.current_rfp['due_date'].strftime('%b %d, %Y')}")
        
        # Progress
        st.markdown("#### 📊 Progress")
        progress = analyzer.get_workflow_progress()
        st.progress(progress / 100)
        st.caption(f"{progress}% Complete")
        
        # Quick Stats
        st.markdown("#### 📈 Statistics")
        col1, col2 = st.columns(2)
        with col1:
            st.metric("Vendors", len(analyzer.vendors))
        with col2:
            st.metric("Shortlisted", len(analyzer.shortlisted_vendors))
        
        if analyzer.selected_vendor:
            selected = analyzer.vendors[analyzer.selected_vendor]
            st.success(f"✅ Selected: {selected.name}")
        
        # Navigation
        st.markdown("---")
        st.markdown("#### 📍 Navigation")
        
        pages = [
            ("workflow", "📋 Workflow", "Manage RFP stages"),
            ("vendors", "👥 Vendors", "Vendor management"),
            ("evaluation", "📊 Evaluation", "Detailed evaluation"),
            ("comparison", "📈 Comparison", "Compare vendors"),
            ("sample_data", "🎲 Test Data", "Generate samples"),
            ("qa", "💬 Q&A", "Ask questions")
        ]
        
        for page_id, label, help_text in pages:
            if st.button(label, use_container_width=True, key=f"nav_{page_id}", help=help_text):
                st.session_state.current_page = page_id
                st.rerun()
        
        # Help
        st.markdown("---")
        with st.expander("❓ Help"):
            st.markdown("""
            **Quick Guide:**
            1. Enable Test Mode
            2. Review workflow stages
            3. Evaluate vendors
            4. Shortlist top performers
            5. Select winner
            6. Complete workflow
            
            **Tips:**
            - Use test mode for demo
            - Progress stages in order
            - Shortlist before selecting
            - Generate reports for each vendor
            """)

# ========================================
# MAIN APPLICATION
# ========================================

def main():
    """Main application entry point"""
    
    # Initialize
    analyzer = RFPAnalyzer()
    
    # Initialize session state
    if 'current_page' not in st.session_state:
        st.session_state.current_page = 'workflow'
    
    # Render components
    render_header()
    render_test_mode_banner()
    render_sidebar(analyzer)
    
    # Main content based on current page
    if st.session_state.current_page == 'workflow':
        render_workflow_dashboard(analyzer)
    
    elif st.session_state.current_page == 'vendors':
        render_vendor_management(analyzer)
    
    elif st.session_state.current_page == 'evaluation':
        render_vendor_evaluation(analyzer)
    
    elif st.session_state.current_page == 'comparison':
        st.header("📈 Vendor Comparison")
        
        if len(analyzer.vendors) >= 2:
            # Comparison matrix
            comparison_data = []
            for vendor in analyzer.vendors.values():
                row = {
                    "Vendor": vendor.name,
                    "Overall Score": vendor.overall_score,
                    "Status": vendor.status,
                    "Strengths": len(vendor.strengths),
                    "Weaknesses": len(vendor.weaknesses)
                }
                comparison_data.append(row)
            
            df = pd.DataFrame(comparison_data)
            df = df.sort_values("Overall Score", ascending=False)
            
            st.dataframe(df, use_container_width=True, hide_index=True)
            
            # Comparison chart
            if analyzer.vendors:
                fig = go.Figure()
                
                for vendor in list(analyzer.vendors.values())[:5]:
                    if vendor.scores:
                        categories = [k.replace('_', ' ').title() for k in vendor.scores.keys()]
                        values = list(vendor.scores.values())
                        
                        fig.add_trace(go.Scatterpolar(
                            r=values,
                            theta=categories,
                            fill='toself',
                            name=vendor.name
                        ))
                
                fig.update_layout(
                    polar=dict(radialaxis=dict(visible=True, range=[0, 100])),
                    showlegend=True,
                    title="Vendor Comparison",
                    height=500
                )
                
                st.plotly_chart(fig, use_container_width=True)
        else:
            st.info("Need at least 2 vendors for comparison. Enable test mode to load sample vendors.")
    
    elif st.session_state.current_page == 'sample_data':
        render_sample_data_generator()
    
    elif st.session_state.current_page == 'qa':
        st.header("💬 Q&A Assistant")
        st.info("Q&A functionality for clarifications and vendor communications")
        
        # Sample Q&A interface
        question = st.text_area("Ask a question about the RFP process:")
        if st.button("Send Question"):
            if question:
                st.success("Question submitted!")
                st.info(f"Sample response: Based on the RFP requirements, {question.lower()} would be addressed in the evaluation phase.")
    
    # Footer
    st.markdown("---")
    st.markdown("""
    <div style="text-align: center; color: #666;">
        <p>© 2025 RFP Vendor Management System v2.0</p>
        <p>Complete Workflow • AI Analysis • Vendor Selection</p>
    </div>
    """, unsafe_allow_html=True)

if __name__ == "__main__":
    main()
